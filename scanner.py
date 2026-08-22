import os
import re
from collections import Counter

import pytesseract
from PIL import Image

# ============================================================
# TESSERACT CONFIGURATION
# ============================================================

TESSERACT_PATH = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

if os.path.exists(TESSERACT_PATH):
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH


# ============================================================
# SEVERITY
# ============================================================

SEVERITY = {
    "Email Address": "High",
    "Phone Number": "High",
    "IP Address": "Medium",
    "Aadhaar Number": "Critical",
    "PAN Number": "Critical",
    "Credit Card": "Critical",
    "Bank Account": "Critical",
    "IFSC Code": "High",
    "Passport Number": "Critical",
    "Password": "Critical",
    "Private Key": "Critical",
    "Cloud Access Key": "Critical",
    "JWT Token": "Critical",
    "UPI ID": "High",
    "GPS Coordinates": "High",
    "MAC Address": "Medium",
    "URL": "Medium",
}


# ============================================================
# REGEX PATTERNS
# ============================================================

PATTERNS = {
    "Email Address": re.compile(
        r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b"
    ),

    "Phone Number": re.compile(
        r"(?<!\d)(?:\+91[\s-]?)?[6-9]\d{9}(?!\d)"
    ),

    "IP Address": re.compile(
        r"\b(?:(?:25[0-5]|2[0-4]\d|1?\d?\d)\.){3}"
        r"(?:25[0-5]|2[0-4]\d|1?\d?\d)\b"
    ),

    "PAN Number": re.compile(
        r"\b[A-Z]{5}[0-9]{4}[A-Z]\b",
        re.IGNORECASE
    ),

    "Credit Card": re.compile(
        r"(?<!\d)(?:\d{4}[\s-]?){3}\d{4}(?!\d)"
    ),

    "Bank Account": re.compile(
        r"(?i)(?:account\s*(?:number|no|#)?\s*[:\-]?\s*)"
        r"\b\d{9,18}\b"
    ),

    "IFSC Code": re.compile(
        r"\b[A-Z]{4}0[A-Z0-9]{6}\b",
        re.IGNORECASE
    ),

    "Passport Number": re.compile(
        r"(?i)(?:passport\s*(?:number|no)?\s*[:\-]?\s*)"
        r"\b[A-Z][0-9]{7}\b"
    ),

    "Password": re.compile(
        r"(?i)(?:password|passwd|pwd|passcode)"
        r"\s*[:=]\s*[^\s,;]+"
    ),

    "Private Key": re.compile(
        r"-----BEGIN (?:RSA |EC |DSA |OPENSSH )?PRIVATE KEY-----"
        r"[\s\S]*?"
        r"-----END (?:RSA |EC |DSA |OPENSSH )?PRIVATE KEY-----",
        re.IGNORECASE
    ),

    "Cloud Access Key": re.compile(
        r"\bAKIA[0-9A-Z]{16}\b"
    ),

    "JWT Token": re.compile(
        r"\beyJ[A-Za-z0-9_-]+\."
        r"[A-Za-z0-9_-]+\."
        r"[A-Za-z0-9_-]+\b"
    ),

    "UPI ID": re.compile(
        r"\b[A-Za-z0-9._-]+@[A-Za-z][A-Za-z0-9._-]{1,30}\b"
    ),

    "GPS Coordinates": re.compile(
        r"(?<!\d)[-+]?(?:[0-8]?\d(?:\.\d+)?|90(?:\.0+)?)"
        r"\s*[,]\s*"
        r"[-+]?(?:1[0-7]\d(?:\.\d+)?|180(?:\.0+)?|"
        r"[0-9]?\d(?:\.\d+)?)"
        r"(?!\d)"
    ),

    "MAC Address": re.compile(
        r"\b(?:[0-9A-Fa-f]{2}[:-]){5}[0-9A-Fa-f]{2}\b"
    ),

    "URL": re.compile(
        r"\bhttps?://[^\s<>\"]+",
        re.IGNORECASE
    ),
}


# ============================================================
# HELPER FUNCTIONS
# ============================================================

def clean_value(value):
    return " ".join(value.strip().split())


def looks_like_credit_card(value):
    digits = re.sub(r"\D", "", value)
    return len(digits) in (13, 14, 15, 16, 17, 18, 19)


def valid_aadhaar_candidate(value):
    digits = re.sub(r"\D", "", value)

    if len(digits) != 12:
        return False

    # Prevent common credit-card fragments from becoming Aadhaar.
    if looks_like_credit_card(value):
        return False

    return True


def extract_aadhaar(text):
    results = []

    pattern = re.compile(
        r"(?<![\d])\d{4}[\s-]?\d{4}[\s-]?\d{4}(?![\d])"
    )

    for match in pattern.finditer(text):
        value = clean_value(match.group())

        # Check surrounding line.
        start = text.rfind("\n", 0, match.start()) + 1
        end = text.find("\n", match.end())

        if end == -1:
            end = len(text)

        line = text[start:end].lower()

        # Don't classify card numbers as Aadhaar.
        if "credit card" in line:
            continue

        if "card number" in line:
            continue

        if "card" in line and "aadhaar" not in line:
            continue

        if valid_aadhaar_candidate(value):
            results.append(value)

    return results


def extract_findings(text):
    findings = []

    # Aadhaar separately to avoid credit-card false positives.
    for value in extract_aadhaar(text):
        findings.append({
            "category": "Aadhaar Number",
            "value": value,
            "severity": SEVERITY["Aadhaar Number"]
        })

    for category, pattern in PATTERNS.items():

        for match in pattern.finditer(text):
            value = clean_value(match.group())

            # Don't treat email as UPI.
            if category == "UPI ID":
                if "." in value.split("@")[0]:
                    continue

                if value.lower().endswith((".com", ".in", ".org", ".net")):
                    continue

            # Don't detect PAN inside other values.
            if category == "PAN Number":
                if len(value) != 10:
                    continue

            # Credit card validation.
            if category == "Credit Card":
                digits = re.sub(r"\D", "", value)

                if not 13 <= len(digits) <= 19:
                    continue

            findings.append({
                "category": category,
                "value": value,
                "severity": SEVERITY[category]
            })

    # Remove duplicates.
    unique = []
    seen = set()

    for finding in findings:
        key = (
            finding["category"],
            finding["value"].lower()
        )

        if key not in seen:
            seen.add(key)
            unique.append(finding)

    return unique


# ============================================================
# RISK SCORE
# ============================================================

SEVERITY_POINTS = {
    "Critical": 20,
    "High": 12,
    "Medium": 6,
    "Low": 2
}


def calculate_risk(findings):
    score = 0

    for finding in findings:
        score += SEVERITY_POINTS.get(
            finding["severity"],
            1
        )

    score = min(score, 100)

    if score >= 75:
        level = "Critical"
    elif score >= 50:
        level = "High"
    elif score >= 25:
        level = "Medium"
    else:
        level = "Low"

    return score, level


# ============================================================
# RISK BREAKDOWN
# ============================================================

def build_risk_breakdown(findings):
    counter = Counter(
        finding["category"]
        for finding in findings
    )

    return dict(counter)


# ============================================================
# RECOMMENDATIONS
# ============================================================

RECOMMENDATIONS = {
    "Email Address":
        "Remove or redact email addresses before sharing.",

    "Phone Number":
        "Remove personal phone numbers before sharing.",

    "IP Address":
        "Consider removing internal or personal IP addresses.",

    "Aadhaar Number":
        "Never publicly share Aadhaar numbers. Mask sensitive digits.",

    "PAN Number":
        "Avoid sharing PAN numbers in public documents.",

    "Credit Card":
        "Mask credit card information before sharing.",

    "Bank Account":
        "Never publicly expose bank account numbers.",

    "IFSC Code":
        "Avoid exposing banking information unnecessarily.",

    "Passport Number":
        "Protect passport numbers from public exposure.",

    "Password":
        "Remove passwords and credentials immediately.",

    "Private Key":
        "Private keys must never be publicly shared.",

    "Cloud Access Key":
        "Revoke and rotate exposed cloud access keys immediately.",

    "JWT Token":
        "Never expose authentication tokens.",

    "UPI ID":
        "Avoid exposing personal UPI IDs in publicly shared documents.",

    "GPS Coordinates":
        "Remove precise GPS coordinates when location privacy is important.",

    "MAC Address":
        "Avoid exposing device MAC addresses.",

    "URL":
        "Check URLs for private or internal information before sharing.",
}


def build_recommendations(findings):
    recommendations = []
    categories = set(
        finding["category"]
        for finding in findings
    )

    for category in RECOMMENDATIONS:
        if category in categories:
            recommendations.append(
                RECOMMENDATIONS[category]
            )

    if not recommendations:
        recommendations.append(
            "No major privacy risks were detected."
        )

    return recommendations


# ============================================================
# IMAGE OCR
# ============================================================

def ocr_image(image):
    try:
        if isinstance(image, str):
            image = Image.open(image)

        if image.mode != "RGB":
            image = image.convert("RGB")

        text = pytesseract.image_to_string(
            image,
            config="--psm 6"
        )

        return text.strip()

    except Exception as error:
        return f"OCR Error: {error}"


def scan_image(image_path):
    text = ocr_image(image_path)

    if text.startswith("OCR Error:"):
        return {
            "success": False,
            "error": text
        }

    return scan_content(
        text,
        file_type=os.path.splitext(image_path)[1].lower()
    )


# ============================================================
# PDF EXTRACTION
# ============================================================

def extract_pdf_text(path):
    text_parts = []

    # First attempt: normal PDF text extraction.
    try:
        from pypdf import PdfReader

        reader = PdfReader(path)

        for page in reader.pages:
            page_text = page.extract_text()

            if page_text:
                text_parts.append(page_text)

    except Exception:
        pass

    extracted_text = "\n".join(text_parts).strip()

    # OCR scanned/image PDF pages.
    try:
        import fitz

        document = fitz.open(path)

        for page_number, page in enumerate(document):

            existing_text = page.get_text("text").strip()

            # OCR if page contains little/no readable text.
            if len(existing_text) < 30:
                pixmap = page.get_pixmap(
                    matrix=fitz.Matrix(2, 2),
                    alpha=False
                )

                image = Image.frombytes(
                    "RGB",
                    [
                        pixmap.width,
                        pixmap.height
                    ],
                    pixmap.samples
                )

                ocr_text = ocr_image(image)

                if ocr_text:
                    text_parts.append(
                        f"\n[OCR Page {page_number + 1}]\n"
                        f"{ocr_text}"
                    )

        document.close()

    except Exception:
        pass

    return "\n".join(text_parts).strip()


# ============================================================
# DOCX EXTRACTION
# ============================================================

def extract_docx_text(path):
    from docx import Document

    document = Document(path)

    parts = []

    # Normal paragraphs.
    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)

    # Tables.
    for table in document.tables:
        for row in table.rows:
            cells = []

            for cell in row.cells:
                if cell.text.strip():
                    cells.append(cell.text.strip())

            if cells:
                parts.append(" | ".join(cells))

    # Embedded images.
    for relationship in document.part.rels.values():

        if "image" not in relationship.reltype:
            continue

        try:
            image_data = relationship.target_part.blob

            from io import BytesIO

            image = Image.open(
                BytesIO(image_data)
            )

            ocr_text = ocr_image(image)

            if ocr_text:
                parts.append(
                    "\n[DOCX IMAGE OCR]\n" + ocr_text
                )

        except Exception:
            continue

    return "\n".join(parts).strip()


# ============================================================
# EXCEL EXTRACTION
# ============================================================

def extract_xlsx_text(path):
    from openpyxl import load_workbook

    workbook = load_workbook(
        path,
        data_only=True
    )

    parts = []

    for worksheet in workbook.worksheets:

        parts.append(
            f"[SHEET: {worksheet.title}]"
        )

        for row in worksheet.iter_rows(
            values_only=True
        ):
            values = []

            for value in row:
                if value is not None:
                    values.append(str(value))

            if values:
                parts.append(
                    " | ".join(values)
                )

    return "\n".join(parts).strip()


# ============================================================
# POWERPOINT EXTRACTION
# ============================================================

def extract_pptx_text(path):
    from pptx import Presentation
    from io import BytesIO

    presentation = Presentation(path)

    parts = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        parts.append(
            f"[SLIDE {slide_number}]"
        )

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                if shape.text.strip():
                    parts.append(
                        shape.text.strip()
                    )

            # OCR images inside PowerPoint.
            if getattr(
                shape,
                "shape_type",
                None
            ) == 13:

                try:
                    image = Image.open(
                        BytesIO(
                            shape.image.blob
                        )
                    )

                    ocr_text = ocr_image(image)

                    if ocr_text:
                        parts.append(
                            "[PPT IMAGE OCR]\n"
                            + ocr_text
                        )

                except Exception:
                    pass

    return "\n".join(parts).strip()


# ============================================================
# FILE TEXT EXTRACTION
# ============================================================

def extract_file_text(path):
    extension = os.path.splitext(
        path
    )[1].lower()

    if extension in [".txt", ".log", ".csv"]:
        with open(
            path,
            "r",
            encoding="utf-8",
            errors="ignore"
        ) as file:
            return file.read()

    if extension == ".pdf":
        return extract_pdf_text(path)

    if extension == ".docx":
        return extract_docx_text(path)

    if extension in [".xlsx", ".xlsm"]:
        return extract_xlsx_text(path)

    if extension == ".pptx":
        return extract_pptx_text(path)

    if extension in [
        ".jpg",
        ".jpeg",
        ".png",
        ".webp",
        ".bmp",
        ".tiff",
        ".tif"
    ]:
        return ocr_image(path)

    raise ValueError(
        f"Unsupported file format: {extension}"
    )


# ============================================================
# MAIN SCANNER
# ============================================================

def scan_content(text, file_type="text"):
    if text is None:
        text = ""

    text = str(text)

    findings = extract_findings(text)

    risk_score, risk_level = calculate_risk(
        findings
    )

    risk_breakdown = build_risk_breakdown(
        findings
    )

    recommendations = build_recommendations(
        findings
    )

    return {
        "total_findings": len(findings),
        "findings": findings,
        "risk_score": risk_score,
        "risk_level": risk_level,
        "risk_breakdown": risk_breakdown,
        "recommendations": recommendations,
        "success": True,
        "extracted_text": text,
        "file_type": file_type
    }


def scan_file(path):
    try:
        extension = os.path.splitext(
            path
        )[1].lower()

        text = extract_file_text(path)

        if not text.strip():
            return {
                "success": False,
                "error":
                    "No readable text was found. "
                    "If this is an image/scanned file, "
                    "make sure Tesseract OCR is installed."
            }

        return scan_content(
            text,
            file_type=extension
        )

    except Exception as error:
        return {
            "success": False,
            "error": str(error)
        }